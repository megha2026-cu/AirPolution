"""
generate_ppt.py — Professional presentation generator
Run: python3 .claude/docs/generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT = os.path.join(os.path.dirname(__file__), 'AirQuality_Presentation.pptx')

NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
BLUE   = RGBColor(0x1A, 0x73, 0xE8)
LBLUE  = RGBColor(0x4A, 0x90, 0xD9)
CYAN   = RGBColor(0x00, 0xB0, 0xD8)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF7, 0xFA)
MGRAY  = RGBColor(0xCC, 0xD6, 0xE0)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
RED    = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(sl, rgb):
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def rect(sl, l, t, w, h, fill, line=None, lw=Pt(0)):
    shp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp


def text(sl, txt, l, t, w, h, size=13, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def multiline(sl, lines, l, t, w, h, size=12, color=DGRAY, spacing=Pt(5)):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color


def chrome(sl, title, subtitle=None, accent=BLUE):
    set_bg(sl, (0xF8, 0xF9, 0xFC))
    rect(sl, 0, 0, 13.33, 1.25, NAVY)
    rect(sl, 0, 1.25, 13.33, 0.06, accent)
    text(sl, title, 0.35, 0.12, 12.6, 0.75, size=24, bold=True, color=WHITE)
    if subtitle:
        text(sl, subtitle, 0.35, 0.88, 12.6, 0.35, size=12,
             color=RGBColor(0xB0, 0xC8, 0xE8), italic=True)
    rect(sl, 0, 7.1, 13.33, 0.4, NAVY)
    text(sl, 'Air Quality Monitoring Dashboard  |  Plag Pro, Noida',
         0.2, 7.12, 12.9, 0.34, size=9,
         color=RGBColor(0x80, 0xA0, 0xC0), align=PP_ALIGN.CENTER)


def bullet_card(sl, title, items, l, t, w, h, accent=BLUE):
    rect(sl, l, t, w, 0.42, accent)
    text(sl, title, l+0.12, t+0.06, w-0.2, 0.32, size=12, bold=True, color=WHITE)
    rect(sl, l, t+0.42, w, h-0.42, WHITE, MGRAY, Pt(0.5))
    multiline(sl, [f'   {i}' for i in items],
              l+0.12, t+0.5, w-0.22, h-0.6, size=11.5, color=DGRAY)


def kpi_box(sl, label, value, sublabel, l, t, w=2.9, h=1.6, accent=BLUE):
    rect(sl, l, t, w, h, WHITE, accent, Pt(1.5))
    rect(sl, l, t, w, 0.34, accent)
    text(sl, label, l+0.1, t+0.05, w-0.18, 0.26, size=10, bold=True, color=WHITE)
    text(sl, value, l+0.1, t+0.4, w-0.18, 0.72, size=26, bold=True,
         color=accent, align=PP_ALIGN.CENTER)
    text(sl, sublabel, l+0.1, t+1.15, w-0.18, 0.3, size=9,
         color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


def slide_table(sl, headers, rows, col_widths, top=1.45):
    all_rows = [headers] + rows
    tbl = sl.shapes.add_table(
        len(all_rows), len(headers),
        Inches(0.3), Inches(top),
        Inches(sum(col_widths)), Inches(len(all_rows) * 0.42)
    ).table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = Inches(w)
    for ri, row in enumerate(all_rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(10)
            run.font.bold = (ri == 0)
            run.font.color.rgb = WHITE if ri == 0 else DGRAY
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = NAVY
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
            else:
                cell.fill.fore_color.rgb = WHITE


# ── SLIDE 1 — TITLE ──────────────────────────────────────────────────────────
s1 = new_slide()
set_bg(s1, (0x1B, 0x2A, 0x4A))
rect(s1, 0, 5.8, 13.33, 0.06, CYAN)
rect(s1, 0, 5.86, 13.33, 1.64, RGBColor(0x12, 0x1E, 0x38))
rect(s1, 1.0, 1.0, 11.33, 3.8, RGBColor(0x22, 0x3A, 0x5E))
rect(s1, 1.0, 1.0, 0.12, 3.8, CYAN)

text(s1, 'Air Quality Monitoring &', 1.25, 1.2, 11.0, 1.1,
     size=32, bold=True, color=WHITE)
text(s1, 'Pollution Trend Visualization Dashboard', 1.25, 2.2, 11.0, 1.1,
     size=26, bold=True, color=CYAN)
text(s1, 'Project Presentation  —  May 2026', 1.25, 3.35, 11.0, 0.5,
     size=15, color=RGBColor(0xB0, 0xC8, 0xE8))

text(s1, 'Plag Pro, Noida', 1.2, 6.08, 6.0, 0.38,
     size=12, color=RGBColor(0x80, 0xA8, 0xD0))
text(s1, 'May 2026', 9.5, 6.08, 3.5, 0.38,
     size=20, bold=True, color=RGBColor(0x60, 0x90, 0xC0), align=PP_ALIGN.RIGHT)

# ── SLIDE 2 — AGENDA ─────────────────────────────────────────────────────────
s2 = new_slide()
chrome(s2, 'Agenda')

agenda = [
    ('01', 'Project Overview & Objectives'),
    ('02', 'System Architecture'),
    ('03', 'Database Design & Schema'),
    ('04', 'ETL Pipeline — Data Ingestion'),
    ('05', 'REST API — Endpoints & Security'),
    ('06', 'Interactive Dashboard & Visualizations'),
    ('07', 'KPIs, Forecasting & Alert System'),
    ('08', 'System Testing Results'),
    ('09', 'Environmental Insights'),
    ('10', 'Deliverables & Conclusion'),
]
for i, (num, title) in enumerate(agenda):
    col, row = i % 2, i // 2
    x = 0.4 + col * 6.5
    y = 1.5 + row * 0.57
    rect(s2, x, y, 0.52, 0.42, BLUE)
    text(s2, num, x, y+0.04, 0.52, 0.36, size=12, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    rect(s2, x+0.52, y, 5.75, 0.42, WHITE, MGRAY, Pt(0.5))
    text(s2, title, x+0.65, y+0.08, 5.5, 0.3, size=12, color=NAVY)

# ── SLIDE 3 — PROJECT OVERVIEW ───────────────────────────────────────────────
s3 = new_slide()
chrome(s3, 'Project Overview & Objectives',
       'Full-stack air quality monitoring solution')

rect(s3, 0.3, 1.45, 8.1, 5.3, WHITE, MGRAY, Pt(0.5))
rect(s3, 0.3, 1.45, 8.1, 0.38, NAVY)
text(s3, 'What we built', 0.45, 1.5, 7.9, 0.3, size=12, bold=True, color=WHITE)
multiline(s3, [
    '  Interactive web dashboard monitoring 10 major Indian cities',
    '  Tracks AQI, PM2.5, PM10, CO, NO2, SO2, O3, Temperature & Humidity',
    '  Full pipeline: simulation  ETL  MySQL  REST API  Chart.js',
    '  7-day AQI forecast with confidence scoring per city',
    '  Automated hazard alerts when AQI exceeds 300',
    '  Filters: city, date range, pollutant type, AQI category',
    '  Responsive layout for desktop and mobile',
    '  27-test automated API test suite — all passing',
], 0.45, 1.9, 7.75, 4.7, size=12, color=DGRAY)

rect(s3, 8.65, 1.45, 4.35, 5.3, WHITE, MGRAY, Pt(0.5))
rect(s3, 8.65, 1.45, 4.35, 0.38, NAVY)
text(s3, 'At a glance', 8.8, 1.5, 4.1, 0.3, size=12, bold=True, color=WHITE)
for i, (val, lbl) in enumerate([
    ('10',    'Cities monitored'),
    ('9',     'Pollutants tracked'),
    ('1,240', 'Data records seeded'),
    ('11',    'API endpoints'),
    ('6',     'Dashboard charts'),
    ('27/27', 'Tests passing'),
    ('14/14', 'Requirements met'),
]):
    y = 1.95 + i * 0.72
    rect(s3, 8.75, y, 1.1, 0.55, BLUE)
    text(s3, val, 8.75, y+0.04, 1.1, 0.48, size=14, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    text(s3, lbl, 9.92, y+0.12, 2.9, 0.34, size=11, color=DGRAY)

# ── SLIDE 4 — SYSTEM ARCHITECTURE ────────────────────────────────────────────
s4 = new_slide()
chrome(s4, 'System Architecture', '3-Tier — Presentation  Application  Data')

for label, color, x, items in [
    ('PRESENTATION TIER', BLUE,  0.3,  ['HTML5 + CSS3', 'Vanilla JavaScript', 'Chart.js (6 charts)', 'Responsive CSS Grid', 'Auto-refresh every 5 min']),
    ('APPLICATION TIER',  GREEN, 4.55, ['Node.js v18 + Express', 'Joi input validation', 'Helmet security headers', 'Rate limiting (100 req/15 min)', 'node-cron scheduled jobs']),
    ('DATA TIER',         NAVY,  8.8,  ['MySQL 8 database', '4 normalized tables', 'Time-series indexed', 'Generated aqi_category col', 'FK constraints + CASCADE']),
]:
    rect(s4, x, 1.45, 4.1, 0.5, color)
    text(s4, label, x+0.1, 1.5, 3.9, 0.4, size=11, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    rect(s4, x, 1.95, 4.1, 4.5, WHITE, color, Pt(1.5))
    multiline(s4, [f'    {i}' for i in items],
              x+0.18, 2.1, 3.8, 4.2, size=12, color=DGRAY)

text(s4, '', 4.2, 3.8, 0.45, 0.5)
text(s4, '>', 4.22, 3.72, 0.4, 0.5, size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
text(s4, '>', 8.45, 3.72, 0.4, 0.5, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── SLIDE 5 — DATABASE SCHEMA ─────────────────────────────────────────────────
s5 = new_slide()
chrome(s5, 'Database Design', 'MySQL 8  —  4 Normalized Tables  —  air_quality_db')

for tname, color, x, y, fields in [
    ('locations',            GREEN,  0.3,  1.45, ['id  PK', 'city', 'state', 'country', 'latitude', 'longitude', 'created_at']),
    ('air_quality_readings', BLUE,   4.5,  1.45, ['id  PK', 'location_id  FK', 'recorded_at', 'aqi', 'pm25 / pm10', 'co / no2 / so2 / o3', 'temperature / humidity', 'aqi_category  (computed)', 'source']),
    ('aqi_forecasts',        PURPLE, 4.5,  4.7,  ['id  PK', 'location_id  FK', 'forecast_date', 'predicted_aqi', 'confidence_pct', 'model_version']),
    ('alerts',               RED,    9.0,  1.45, ['id  PK', 'location_id  FK', 'triggered_at', 'aqi_threshold', 'actual_aqi', 'message', 'resolved_at']),
]:
    h = len(fields) * 0.37 + 0.5
    rect(s5, x, y, 3.9, 0.42, color)
    text(s5, tname, x+0.12, y+0.07, 3.7, 0.3, size=12, bold=True, color=WHITE)
    rect(s5, x, y+0.42, 3.9, h-0.42, WHITE, color, Pt(1.2))
    multiline(s5, fields, x+0.16, y+0.5, 3.68, h-0.55, size=11, color=DGRAY, spacing=Pt(3))

text(s5, '1  N', 4.05, 2.65, 0.55, 0.28, size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
text(s5, '1  N', 4.05, 5.55, 0.55, 0.28, size=9, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
text(s5, '1  N', 8.6,  2.65, 0.55, 0.28, size=9, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ── SLIDE 6 — ETL PIPELINE ───────────────────────────────────────────────────
s6 = new_slide()
chrome(s6, 'ETL Pipeline', 'data/etl_import.js  —  CSV to MySQL with data cleaning')

for label, color, x, detail in [
    ('1\nEXTRACT',   BLUE,  0.3,  ['Read CSV as stream', 'Auto-detect columns', 'Any CPCB/WHO format', 'Memory efficient']),
    ('2\nTRANSFORM', GREEN, 3.55, ['Validate dates', 'Detect outliers', 'Clamp to valid range', 'Handle missing  NULL', 'Normalise units']),
    ('3\nLOAD',      NAVY,  6.8,  ['Batch 200 rows/query', 'Auto-create cities', 'INSERT IGNORE dupes', 'Dry-run mode']),
]:
    rect(s6, x, 1.45, 3.0, 5.3, color)
    text(s6, label, x+0.1, 1.55, 2.8, 0.75, size=14, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    rect(s6, x+0.12, 2.38, 2.76, 4.25, RGBColor(0xFF, 0xFF, 0xFF))
    multiline(s6, [f'    {d}' for d in detail],
              x+0.18, 2.48, 2.65, 4.0, size=12, color=DGRAY)

text(s6, '>', 3.2,  3.55, 0.42, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s6, '>', 6.45, 3.55, 0.42, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s6, 10.1, 1.45, 2.9, 5.3, WHITE, MGRAY, Pt(0.5))
rect(s6, 10.1, 1.45, 2.9, 0.38, RGBColor(0x37, 0x47, 0x5A))
text(s6, 'Summary Output', 10.2, 1.5, 2.7, 0.3, size=11, bold=True, color=WHITE)
multiline(s6, [
    'Lines processed', 'Rows inserted', 'Rows skipped',
    'Outliers found & clamped', '',
    'Usage:', 'node etl_import.js', '--file=data.csv', '[--dry-run]',
], 10.22, 1.9, 2.65, 4.7, size=11, color=DGRAY)

# ── SLIDE 7 — API ENDPOINTS ───────────────────────────────────────────────────
s7 = new_slide()
chrome(s7, 'REST API — Endpoints & Security',
       'Node.js v18 + Express  |  Port 3000  |  11 Endpoints')

slide_table(s7,
    ['Method', 'Endpoint', 'Description'],
    [
        ['GET',  '/api/health',                       'Server health check'],
        ['GET',  '/api/locations',                    'List all monitored cities'],
        ['POST', '/api/locations',                    'Register a new city'],
        ['GET',  '/api/readings',                     'Raw readings filtered by city, date, pollutant'],
        ['GET',  '/api/readings/latest',              'Most recent reading per city'],
        ['GET',  '/api/readings/daily',               'Daily avg / max / min AQI summary'],
        ['GET',  '/api/readings/hourly',              'Hourly average AQI'],
        ['GET',  '/api/readings/peak-hours',          'Peak pollution hour breakdown'],
        ['GET',  '/api/readings/category-breakdown',  'AQI category % distribution'],
        ['GET',  '/api/forecast',                     '7-day AQI forecast per city'],
        ['GET',  '/api/alerts',                       'All active hazard alerts'],
    ],
    [1.8, 5.0, 5.9], top=1.45)

rect(s7, 0.3, 6.35, 12.7, 0.65, RGBColor(0xE8, 0xF0, 0xFE), BLUE, Pt(0.5))
text(s7, 'Security:  Helmet headers   Joi validation   Rate limiting (100 req/IP/15 min)   Parameterized SQL   .env credentials',
     0.5, 6.43, 12.3, 0.44, size=11, color=NAVY)

# ── SLIDE 8 — DASHBOARD & CHARTS ─────────────────────────────────────────────
s8 = new_slide()
chrome(s8, 'Interactive Dashboard & Visualizations',
       'Chart.js  |  Vanilla JS  |  Responsive CSS Grid')

for i, (title, color, ctype, desc) in enumerate([
    ('AQI Trend',         BLUE,   'Line chart',      'Historical AQI over selected date range'),
    ('Daily Summary',     GREEN,  'Bar + Line',      'Daily avg AQI with trend overlay'),
    ('Hourly Pattern',    NAVY,   'Bar chart',       'Average AQI by hour of day (0-23)'),
    ('Peak Hours',        ORANGE, 'Horizontal Bar',  'Top pollution hours ranked by avg AQI'),
    ('Category Breakdown',PURPLE, 'Doughnut chart',  'Good / Moderate / Poor / Severe %'),
    ('7-Day Forecast',    CYAN,   'Line + bands',    'Predicted AQI with confidence range'),
]):
    col, row = i % 3, i // 3
    x = 0.3 + col * 4.35
    y = 1.48 + row * 2.82
    rect(s8, x, y, 4.1, 2.58, WHITE, color, Pt(1.5))
    rect(s8, x, y, 4.1, 0.4, color)
    text(s8, title, x+0.12, y+0.07, 3.88, 0.3, size=12, bold=True, color=WHITE)
    rect(s8, x+0.12, y+0.5, 1.55, 0.28, RGBColor(0xE8, 0xF0, 0xFE))
    text(s8, ctype, x+0.16, y+0.52, 1.5, 0.24, size=9, color=BLUE, bold=True)
    text(s8, desc, x+0.12, y+0.9, 3.85, 1.45, size=11, color=DGRAY, wrap=True)

rect(s8, 0.3, 7.05, 12.7, 0.32, RGBColor(0xF0, 0xF4, 0xF8))
text(s8, 'Filters:  City    Date From/To    Pollutant (AQI / PM2.5 / PM10 / CO / NO2 / SO2 / O3)    Auto-refresh every 5 minutes',
     0.5, 7.08, 12.3, 0.26, size=10, color=NAVY)

# ── SLIDE 9 — KPIs, FORECASTING & ALERTS ─────────────────────────────────────
s9 = new_slide()
chrome(s9, 'KPIs, Forecasting & Alert System',
       'Monitoring, prediction and automated hazard detection')

kpi_box(s9, 'Average AQI',    'AVG(aqi)',  'Over selected date range', 0.3,  1.5, accent=BLUE)
kpi_box(s9, 'Peak AQI',       'MAX(aqi)',  'Highest reading in period', 3.45, 1.5, accent=NAVY)
kpi_box(s9, 'Safe Days',      '<= 100',   'Daily avg AQI <= 100',      6.6,  1.5, accent=GREEN)
kpi_box(s9, 'Hazardous Days', '> 300',    'Daily avg AQI > 300',       9.75, 1.5, accent=RED)

bullet_card(s9, '7-Day AQI Forecast', [
    '14-day moving average as baseline',
    '+/-10 AQI random jitter for realism',
    'Confidence: 90% on day 1, 55% on day 7',
    'Regenerated daily at 01:00 via cron',
    'Stored in aqi_forecasts table',
], 0.3, 3.4, 6.15, 3.3, accent=BLUE)

bullet_card(s9, 'Automated Hazard Alert System', [
    'Trigger: AQI > 300 in last 1 hour',
    'Checked every 30 minutes via node-cron',
    'No repeat alert within 3 hours per city',
    'Stored in alerts table with timestamp',
    'Pulsing red badge shown on dashboard',
], 6.75, 3.4, 6.25, 3.3, accent=RED)

# ── SLIDE 10 — TEST RESULTS ───────────────────────────────────────────────────
s10 = new_slide()
chrome(s10, 'System Testing Results',
       'tests/test_api.js  |  27 / 27 Assertions Passing')

for i, (label, score) in enumerate([
    ('Health Check',        '2 / 2'),
    ('Locations API',       '4 / 4'),
    ('Latest Readings',     '3 / 3'),
    ('Filtered Readings',   '2 / 2'),
    ('Daily Summary',       '2 / 2'),
    ('Hourly Average',      '2 / 2'),
    ('Peak Hours',          '2 / 2'),
    ('Category Breakdown',  '2 / 2'),
    ('Forecast',            '3 / 3'),
    ('Alerts',              '2 / 2'),
    ('Input Validation',    '2 / 2'),
    ('Real-Time Stream',    '1 / 1'),
]):
    col, row = i % 4, i // 4
    x = 0.3 + col * 3.25
    y = 1.5 + row * 1.62
    rect(s10, x, y, 3.05, 1.45, WHITE, GREEN, Pt(1.2))
    rect(s10, x, y, 3.05, 0.38, GREEN)
    text(s10, label, x+0.1, y+0.07, 2.85, 0.28, size=10, bold=True, color=WHITE)
    text(s10, score, x+0.1, y+0.5, 2.85, 0.72, size=22, bold=True,
         color=GREEN, align=PP_ALIGN.CENTER)

rect(s10, 0.3, 6.5, 12.7, 0.58, GREEN)
text(s10, 'TOTAL:  27 / 27 Tests Passed     Zero failures     Includes simulated real-time sensor stream',
     0.3, 6.58, 12.7, 0.42, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDE 11 — ENVIRONMENTAL INSIGHTS ────────────────────────────────────────
s11 = new_slide()
chrome(s11, 'Environmental Insights',
       'Key findings from the dashboard analytics')

for i, (color, title, body_txt) in enumerate([
    (RED,    'Most Polluted Cities',   'Delhi, Noida and Gurugram show AQI in "Poor" to "Very Poor" range (201-400) consistently across the monitoring period.'),
    (GREEN,  'Cleanest Cities',        'Bengaluru and Chennai average in the "Moderate" category (101-200) with significantly more safe days per month.'),
    (ORANGE, 'Peak Pollution Hours',   'AQI spikes observed during 6-9 AM and 7-10 PM correlating directly with vehicle traffic rush hours.'),
    (BLUE,   'Key Pollutant',          'PM2.5 is the dominant contributor to high AQI across all cities, driven by vehicular and industrial emissions.'),
    (PURPLE, 'Seasonal Pattern',       'Winter months record higher AQI values due to temperature inversion trapping pollutants near ground level.'),
    (NAVY,   'Hazardous Events',       'AQI > 300 events are concentrated in northern cities (Delhi, Gurugram, Noida) during November to January.'),
]):
    col, row = i % 2, i // 2
    x = 0.3 + col * 6.55
    y = 1.48 + row * 1.9
    rect(s11, x, y, 6.3, 1.74, WHITE, color, Pt(1.2))
    rect(s11, x, y, 0.22, 1.74, color)
    text(s11, title, x+0.35, y+0.1, 5.8, 0.35, size=12, bold=True, color=color)
    text(s11, body_txt, x+0.35, y+0.5, 5.8, 1.1, size=11, color=DGRAY, wrap=True)

# ── SLIDE 12 — DELIVERABLES ───────────────────────────────────────────────────
s12 = new_slide()
chrome(s12, 'Deliverables & Conclusion', 'All 14 requirements fulfilled')

for i, d in enumerate([
    'Cleaned dataset CSV  —  data/sample/air_quality_sample.csv',
    'ETL Pipeline  —  data/etl_import.js',
    'ER Diagram  —  docs/er_diagram.html',
    'SQL Schema  —  database/migrations/001_create_tables.sql',
    'REST API (11 endpoints)  —  backend/',
    'Interactive Dashboard (6 charts)  —  frontend/index.html',
    'Forecasting Model  —  forecastController.js',
    'Alert System  —  alertController.js',
    'Test Suite 27/27  —  tests/test_api.js',
    'PDF Report  —  docs/AirQuality_Project_Report.pdf',
    'PPT Presentation  —  docs/AirQuality_Presentation.pptx',
    'System Architecture Doc  —  docs/SYSTEM_ARCHITECTURE.md',
]):
    col, row = i % 2, i // 2
    x = 0.3 + col * 6.5
    y = 1.5 + row * 0.57
    rect(s12, x, y, 0.38, 0.42, GREEN)
    text(s12, 'v', x, y+0.05, 0.38, 0.36, size=13, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    rect(s12, x+0.38, y, 6.0, 0.42, WHITE, MGRAY, Pt(0.3))
    text(s12, d, x+0.5, y+0.09, 5.75, 0.28, size=10, color=NAVY)

rect(s12, 0.3, 6.42, 12.7, 0.58, NAVY)
text(s12,
     'Complete production-quality platform covering all 14 requirements with full test coverage.',
     0.5, 6.5, 12.3, 0.44, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDE 13 — THANK YOU ──────────────────────────────────────────────────────
s13 = new_slide()
set_bg(s13, (0x1B, 0x2A, 0x4A))
rect(s13, 0, 0, 13.33, 7.5, NAVY)

# "Thank You" large centered text
text(s13, 'Thank You', 0, 1.6, 13.33, 2.0,
     size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Single thin cyan line as divider
rect(s13, 0, 4.0, 13.33, 0.04, CYAN)

# Only project subtitle below line — no group, no org, no stats
text(s13, 'Air Quality Monitoring & Pollution Trend Visualization Dashboard',
     0.5, 4.2, 12.33, 0.6,
     size=15, color=CYAN, align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f'PPT generated: {OUTPUT}  ({len(prs.slides)} slides)')
